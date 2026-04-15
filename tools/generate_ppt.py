#!/usr/bin/env python3
"""
Nutonium – Hyperlocal Social Commerce Platform
Professional PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────────────
DEEP_NAVY    = RGBColor(0x0B, 0x0F, 0x1A)
DARK_BG      = RGBColor(0x10, 0x15, 0x24)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xA0, 0xAE, 0xC0)
MID_GREY     = RGBColor(0x64, 0x74, 0x8B)
CARD_BG      = RGBColor(0x1A, 0x1F, 0x2E)
GREEN        = RGBColor(0x10, 0xB9, 0x81)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT   = RGBColor(0xEF, 0x44, 0x44)
SOFT_WHITE   = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────────────────────

def _solid_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _gradient_bg(slide):
    """Solid dark background (gradient fills need XML hacking; keep it clean)."""
    _solid_bg(slide, DEEP_NAVY)


def _add_rect(slide, left, top, width, height, fill_colour, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.04
    return shape


def _add_line(slide, left, top, width, colour, thickness=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def _add_circle(slide, left, top, size, fill_colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18, colour=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     colour=SOFT_WHITE, bullet_colour=ACCENT_CYAN, spacing=Pt(10)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
        # Add bullet using Unicode character
        run = p.runs[0]
        run.text = f"▸  {item}"
    return txBox


def _slide_number(slide, num, total):
    _add_textbox(slide, Inches(12.2), Inches(6.9), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, colour=MID_GREY,
                 alignment=PP_ALIGN.RIGHT)


def _section_header_decor(slide):
    """Add decorative accent line at top of slide."""
    _add_line(slide, Inches(0.8), Inches(0.6), Inches(2.5), ACCENT_CYAN, Pt(4))


def _add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                   accent=ACCENT_CYAN):
    card = _add_rect(slide, left, top, width, height, CARD_BG)
    # Icon circle
    _add_circle(slide, left + Inches(0.25), top + Inches(0.25), Inches(0.55), accent)
    _add_textbox(slide, left + Inches(0.25), top + Inches(0.25), Inches(0.55), Inches(0.55),
                 icon_text, font_size=18, colour=DEEP_NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name='Segoe UI Emoji')
    # Title
    _add_textbox(slide, left + Inches(0.2), top + Inches(1.0), width - Inches(0.4), Inches(0.4),
                 title, font_size=15, colour=WHITE, bold=True)
    # Description
    _add_textbox(slide, left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), height - Inches(1.6),
                 desc, font_size=12, colour=LIGHT_GREY)
    return card


TOTAL_SLIDES = 16

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_gradient_bg(slide)

# Decorative circles (background flourish)
_add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(5), RGBColor(0x00, 0xD4, 0xFF))
# Make it semi-transparent by overlaying dark rect
overlay = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5), Inches(5), Inches(5))
overlay.fill.solid()
overlay.fill.fore_color.rgb = DEEP_NAVY
overlay.fill.fore_color.brightness = 0.0
overlay.line.fill.background()

_add_circle(slide, Inches(10), Inches(4.5), Inches(5), RGBColor(0x8B, 0x5C, 0xF6))
overlay2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(5), Inches(5))
overlay2.fill.solid()
overlay2.fill.fore_color.rgb = DEEP_NAVY
overlay2.line.fill.background()

# Accent bar top
_add_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT_CYAN, Pt(4))

# Main title
_add_textbox(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(1.2),
             "NUTONIUM", font_size=56, colour=ACCENT_CYAN, bold=True, font_name='Calibri')

_add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.9),
             "Hyperlocal Social Commerce Platform", font_size=32, colour=WHITE,
             bold=True, font_name='Calibri')

_add_line(slide, Inches(1.5), Inches(3.4), Inches(3), ACCENT_CYAN, Pt(3))

_add_textbox(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "Connecting Retail Shops, Suppliers & Customers through\nReal-time Inventory, Preorders & Live Product Streaming",
             font_size=16, colour=LIGHT_GREY, font_name='Calibri')

# Team box
team_card = _add_rect(slide, Inches(1.5), Inches(4.9), Inches(5), Inches(2.0), CARD_BG)
_add_textbox(slide, Inches(1.7), Inches(5.0), Inches(4.6), Inches(0.35),
             "PRESENTED BY", font_size=11, colour=ACCENT_CYAN, bold=True)
_add_textbox(slide, Inches(1.7), Inches(5.35), Inches(4.6), Inches(1.4),
             "Abhishek Bruno  ·  Muhammad Aslam\nGokul SR  ·  Rijin Reji",
             font_size=14, colour=SOFT_WHITE)

guide_card = _add_rect(slide, Inches(7.0), Inches(4.9), Inches(4.8), Inches(2.0), CARD_BG)
_add_textbox(slide, Inches(7.2), Inches(5.0), Inches(4.4), Inches(0.35),
             "PROJECT GUIDANCE", font_size=11, colour=ACCENT_CYAN, bold=True)
_add_textbox(slide, Inches(7.2), Inches(5.35), Inches(4.4), Inches(1.4),
             "Project Guide: Ms Aswathy\nProject Coordinator: Dr Anoop Sreekumar",
             font_size=14, colour=SOFT_WHITE)

_slide_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – ABSTRACT
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "ABSTRACT", font_size=36, colour=WHITE, bold=True)

abstract_items = [
    "Nutonium is a Hyperlocal Social Commerce Platform designed to bridge the gap between offline retail and digital commerce.",
    "Connects Retailers, Suppliers, and Customers in a unified ecosystem with real-time data sharing.",
    "Provides Real-time Inventory Visibility – customers can check product availability at nearby shops before visiting.",
    "Enables Preorders & Live Product Streaming – bringing the in-store experience to digital screens.",
    "Supports B2B Supply Collaboration – suppliers gain demand awareness to optimize supply chains.",
    "Improves Shopping Convenience – reducing wasted trips and stock uncertainty for everyday consumers.",
]

y = Inches(1.8)
for item in abstract_items:
    card = _add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), CARD_BG)
    _add_circle(slide, Inches(1.0), y + Inches(0.15), Inches(0.35), ACCENT_CYAN)
    _add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(0.35), Inches(0.35),
                 "✓", font_size=14, colour=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.55), y + Inches(0.1), Inches(10.5), Inches(0.55),
                 item, font_size=14, colour=SOFT_WHITE)
    y += Inches(0.82)

_slide_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "INTRODUCTION", font_size=36, colour=WHITE, bold=True)

_add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
             "The retail landscape faces a critical disconnect between physical stores and the digital world.",
             font_size=16, colour=LIGHT_GREY)

# Problem cards
problems = [
    ("🏪", "Offline Retail Gap", "Local retail shops lack digital integration, missing out on a huge online-savvy customer base.", ACCENT_BLUE),
    ("❓", "Stock Uncertainty", "Customers face uncertainty about product availability, leading to wasted trips and frustration.", ACCENT_PURPLE),
    ("📊", "No Inventory Tools", "Retailers lack affordable, real-time inventory management and visibility tools.", ACCENT_CYAN),
    ("📦", "Supply Blindness", "Suppliers have no awareness of real-time demand at the local retail level.", ORANGE),
    ("🔗", "Nutonium Solution", "Nutonium bridges offline and online retail with a hyperlocal-first social commerce platform.", GREEN),
]

x = Inches(0.5)
for icon, title, desc, accent in problems:
    _add_icon_card(slide, x, Inches(2.5), Inches(2.3), Inches(3.8), icon, title, desc, accent)
    x += Inches(2.5)

_slide_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "OBJECTIVES", font_size=36, colour=WHITE, bold=True)

objectives = [
    ("01", "Digitalize Offline Retail", "Transform traditional brick-and-mortar shops into digitally-enabled storefronts with product catalogs, real-time stock info, and online presence.", ACCENT_CYAN),
    ("02", "Real-time Stock Updates", "Provide instant inventory visibility so customers always know what's available nearby before they leave home.", ACCENT_BLUE),
    ("03", "Enable Preorder System", "Allow customers to reserve products at local shops, guaranteeing availability and reducing missed sales.", ACCENT_PURPLE),
    ("04", "Improve B2B Communication", "Create a direct channel between suppliers and retailers for demand-driven restocking and collaboration.", ORANGE),
    ("05", "Enhance Customer Experience", "Deliver a seamless, modern shopping experience that combines discovery, social commerce, and local retail.", GREEN),
]

y = Inches(1.8)
for num, title, desc, accent in objectives:
    # Number circle
    _add_circle(slide, Inches(0.8), y + Inches(0.05), Inches(0.6), accent)
    _add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(0.6), Inches(0.6),
                 num, font_size=18, colour=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    _add_textbox(slide, Inches(1.6), y, Inches(3), Inches(0.4),
                 title, font_size=18, colour=WHITE, bold=True)
    # Description
    _add_textbox(slide, Inches(1.6), y + Inches(0.4), Inches(10.5), Inches(0.55),
                 desc, font_size=13, colour=LIGHT_GREY)
    y += Inches(1.05)

_slide_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – EXISTING SYSTEMS & LIMITATIONS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
             "EXISTING SYSTEMS & LIMITATIONS", font_size=36, colour=WHITE, bold=True)

# Comparison cards
systems = [
    ("Traditional Retail", "Offline-only model", [
        "No digital presence",
        "Zero inventory tracking",
        "Limited customer reach",
        "No data-driven decisions"
    ], RED_ACCENT),
    ("E-Commerce Giants", "Amazon, Flipkart, etc.", [
        "No hyperlocal focus",
        "No real-time local stock",
        "No direct B2B interaction",
        "Long delivery timelines"
    ], ORANGE),
    ("Nutonium", "Hyperlocal Social Commerce", [
        "Hyperlocal-first approach",
        "Real-time inventory sync",
        "Integrated B2B channel",
        "Instant local fulfillment"
    ], GREEN),
]

x = Inches(0.5)
for title, subtitle, items, accent in systems:
    card = _add_rect(slide, x, Inches(1.9), Inches(3.9), Inches(4.8), CARD_BG)
    _add_line(slide, x, Inches(1.9), Inches(3.9), accent, Pt(4))
    _add_textbox(slide, x + Inches(0.3), Inches(2.15), Inches(3.3), Inches(0.4),
                 title, font_size=22, colour=WHITE, bold=True)
    _add_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.35),
                 subtitle, font_size=13, colour=LIGHT_GREY)

    y = Inches(3.2)
    for item in items:
        marker = "✗" if accent != GREEN else "✓"
        marker_col = accent
        _add_textbox(slide, x + Inches(0.3), y, Inches(0.3), Inches(0.35),
                     marker, font_size=14, colour=marker_col, bold=True)
        _add_textbox(slide, x + Inches(0.65), y, Inches(3.0), Inches(0.35),
                     item, font_size=14, colour=SOFT_WHITE)
        y += Inches(0.45)

    x += Inches(4.2)

_slide_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – LITERATURE SURVEY (Part 1)
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
             "LITERATURE SURVEY", font_size=36, colour=WHITE, bold=True)

lit_data_1 = [
    ("Hyperlocal E-Commerce: Trends and Opportunities", "V. Kumar, R. Rajan", "IEEE Xplore",
     "Explores hyperlocal e-commerce connecting nearby buyers and sellers.",
     "Lacks real-time inventory sync and supplier collaboration."),
    ("A Study on Social Commerce Platforms", "H. Wang, P. Zhang", "Springer",
     "Examines social interaction's role in online purchasing behavior.",
     "No hyperlocal retail integration or inventory management."),
    ("Real-Time Inventory Management in Retail", "S. Patel, M. Shah", "IJCA",
     "Highlights real-time inventory tracking for operational efficiency.",
     "Limited to inventory without customer/supplier integration."),
    ("B2B Supply Chain Using Digital Platforms", "A. Singh, R. Verma", "IEEE Xplore",
     "Discusses platforms enhancing supplier-retailer collaboration.",
     "No end-user interaction or social commerce capabilities."),
]

# Table header
header_y = Inches(1.7)
_add_rect(slide, Inches(0.5), header_y, Inches(12.3), Inches(0.45), ACCENT_BLUE)
cols = [Inches(0.6), Inches(4.1), Inches(5.9), Inches(8.1), Inches(10.3)]
headers = ["Title", "Authors / Source", "Description", "Limitations"]
widths_h = [Inches(3.4), Inches(1.7), Inches(2.1), Inches(2.1)]
for i, h in enumerate(headers):
    _add_textbox(slide, cols[i], header_y + Inches(0.05), widths_h[i], Inches(0.35),
                 h, font_size=12, colour=WHITE, bold=True)

y = header_y + Inches(0.55)
for title, authors, source, desc, limitation in lit_data_1:
    bg_col = CARD_BG if lit_data_1.index((title, authors, source, desc, limitation)) % 2 == 0 else RGBColor(0x15, 0x1A, 0x29)
    _add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1), bg_col)
    _add_textbox(slide, cols[0], y + Inches(0.08), Inches(3.4), Inches(0.95),
                 title, font_size=11, colour=ACCENT_CYAN, bold=True)
    _add_textbox(slide, cols[1], y + Inches(0.08), Inches(1.7), Inches(0.95),
                 f"{authors}\n({source})", font_size=10, colour=SOFT_WHITE)
    _add_textbox(slide, cols[2], y + Inches(0.08), Inches(2.1), Inches(0.95),
                 desc, font_size=10, colour=LIGHT_GREY)
    _add_textbox(slide, cols[3], y + Inches(0.08), Inches(2.1), Inches(0.95),
                 limitation, font_size=10, colour=ORANGE)
    y += Inches(1.2)

_slide_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – LITERATURE SURVEY (Part 2)
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
             "LITERATURE SURVEY (CONTD.)", font_size=36, colour=WHITE, bold=True)

lit_data_2 = [
    ("Web-Based Retail Management Systems", "J. Anderson, L. Thomas", "Google Scholar",
     "Web systems for managing retail operations and customer interactions.",
     "Lacks hyperlocal discovery and real-time sync."),
    ("Live Streaming Commerce: Future of Shopping", "Y. Chen, X. Liu", "Elsevier",
     "Live product streaming to enhance engagement and boost sales.",
     "Focuses on global markets, not hyperlocal retail."),
    ("Location-Based Services in Mobile Commerce", "R. Gupta, N. Kaur", "IJERT",
     "GPS and location tech for personalized, geographically relevant services.",
     "No inventory tracking or supplier collaboration."),
    ("Integrated E-Commerce & Supply Chain", "K. Lee, D. Kim", "IEEE Xplore",
     "Frameworks combining e-commerce with supply chain management.",
     "Lacks social commerce and hyperlocal retail focus."),
]

header_y = Inches(1.7)
_add_rect(slide, Inches(0.5), header_y, Inches(12.3), Inches(0.45), ACCENT_BLUE)
for i, h in enumerate(headers):
    _add_textbox(slide, cols[i], header_y + Inches(0.05), widths_h[i], Inches(0.35),
                 h, font_size=12, colour=WHITE, bold=True)

y = header_y + Inches(0.55)
for idx, (title, authors, source, desc, limitation) in enumerate(lit_data_2):
    bg_col = CARD_BG if idx % 2 == 0 else RGBColor(0x15, 0x1A, 0x29)
    _add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1), bg_col)
    _add_textbox(slide, cols[0], y + Inches(0.08), Inches(3.4), Inches(0.95),
                 title, font_size=11, colour=ACCENT_CYAN, bold=True)
    _add_textbox(slide, cols[1], y + Inches(0.08), Inches(1.7), Inches(0.95),
                 f"{authors}\n({source})", font_size=10, colour=SOFT_WHITE)
    _add_textbox(slide, cols[2], y + Inches(0.08), Inches(2.1), Inches(0.95),
                 desc, font_size=10, colour=LIGHT_GREY)
    _add_textbox(slide, cols[3], y + Inches(0.08), Inches(2.1), Inches(0.95),
                 limitation, font_size=10, colour=ORANGE)
    y += Inches(1.2)

_slide_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "PROBLEM STATEMENT", font_size=36, colour=WHITE, bold=True)

problems_stmt = [
    ("No Real-time Local Availability",
     "There is no unified platform that allows customers to check product availability at nearby retail shops in real time, resulting in wasted trips and missed sales.",
     "🔍", ACCENT_CYAN),
    ("Inefficient Digital Inventory",
     "Small and medium retailers lack affordable tools to digitally manage and broadcast their inventory, limiting their reach and competitiveness against e-commerce giants.",
     "📋", ACCENT_BLUE),
    ("Supplier Demand Blindness",
     "Suppliers and wholesalers lack visibility into real-time demand at the retail level, leading to overstocking, understocking, and broken supply chains.",
     "📦", ACCENT_PURPLE),
]

y = Inches(1.9)
for title, desc, icon, accent in problems_stmt:
    card = _add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.3), CARD_BG)
    _add_line(slide, Inches(0.8), y, Inches(0.03), accent, Inches(1.3))
    _add_textbox(slide, Inches(1.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.8), y + Inches(0.15), Inches(10), Inches(0.4),
                 title, font_size=20, colour=WHITE, bold=True)
    _add_textbox(slide, Inches(1.8), y + Inches(0.6), Inches(10), Inches(0.6),
                 desc, font_size=13, colour=LIGHT_GREY)
    y += Inches(1.5)

# Conclusion banner
_add_rect(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.7), ACCENT_CYAN)
_add_textbox(slide, Inches(2.5), Inches(6.25), Inches(8.3), Inches(0.6),
             "→  Need for a Unified Hyperlocal Commerce Platform",
             font_size=20, colour=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

_slide_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "SYSTEM ARCHITECTURE", font_size=36, colour=WHITE, bold=True)

_add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "A three-tier architecture connecting users, business logic, and data storage through Firebase services.",
             font_size=15, colour=LIGHT_GREY)

# Presentation layer
_add_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.3), CARD_BG)
_add_line(slide, Inches(0.5), Inches(2.3), Inches(12.3), ACCENT_CYAN, Pt(3))
_add_textbox(slide, Inches(0.7), Inches(2.45), Inches(3), Inches(0.35),
             "PRESENTATION LAYER", font_size=14, colour=ACCENT_CYAN, bold=True)
_add_textbox(slide, Inches(0.7), Inches(2.85), Inches(11.8), Inches(0.6),
             "Flutter Mobile App  ·  Social Feed  ·  Map View  ·  Cart & Checkout  ·  Profile Management  ·  Camera  ·  Trade Code Scanner",
             font_size=13, colour=SOFT_WHITE)

# Arrow
_add_textbox(slide, Inches(5.8), Inches(3.6), Inches(2), Inches(0.5),
             "▼", font_size=28, colour=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Business logic layer
_add_rect(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(1.3), CARD_BG)
_add_line(slide, Inches(0.5), Inches(3.95), Inches(12.3), ACCENT_BLUE, Pt(3))
_add_textbox(slide, Inches(0.7), Inches(4.1), Inches(3), Inches(0.35),
             "BUSINESS LOGIC LAYER", font_size=14, colour=ACCENT_BLUE, bold=True)
_add_textbox(slide, Inches(0.7), Inches(4.5), Inches(11.8), Inches(0.6),
             "BLoC State Management  ·  Auth (Email OTP, Google, Phone)  ·  Marketplace Service  ·  Social Service  ·  Firestore CRUD  ·  Role-Based Access",
             font_size=13, colour=SOFT_WHITE)

# Arrow
_add_textbox(slide, Inches(5.8), Inches(5.25), Inches(2), Inches(0.5),
             "▼", font_size=28, colour=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Data layer
_add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3), CARD_BG)
_add_line(slide, Inches(0.5), Inches(5.6), Inches(12.3), ACCENT_PURPLE, Pt(3))
_add_textbox(slide, Inches(0.7), Inches(5.75), Inches(3), Inches(0.35),
             "DATA & SERVICES LAYER", font_size=14, colour=ACCENT_PURPLE, bold=True)
_add_textbox(slide, Inches(0.7), Inches(6.15), Inches(11.8), Inches(0.6),
             "Firebase Auth  ·  Cloud Firestore (Users, Retailers, Wholesalers, Posts)  ·  Firebase Storage  ·  REST APIs  ·  SMTP Webhooks",
             font_size=13, colour=SOFT_WHITE)

_slide_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – MODULES
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "SYSTEM MODULES", font_size=36, colour=WHITE, bold=True)

modules = [
    ("👤", "User / Customer Module", [
        "Browse social feed & discover shops",
        "Search products across nearby retailers",
        "View shop inventory & locations on map",
        "Add items to cart & preorder",
        "Trade code handoff for quick adds",
    ], ACCENT_CYAN),
    ("🏪", "Retailer Module", [
        "Digital shop profile setup",
        "Real-time inventory management",
        "Accept and manage orders",
        "Post products to social feed",
        "Camera-based product listing",
    ], ACCENT_BLUE),
    ("📦", "Supplier / Wholesaler Module", [
        "Wholesaler profile & catalog setup",
        "Identify retailer stock shortages",
        "Send supply proposals & offers",
        "B2B communication channel",
        "Demand-driven restocking",
    ], ACCENT_PURPLE),
]

x = Inches(0.3)
for icon, title, items, accent in modules:
    card = _add_rect(slide, x, Inches(1.8), Inches(4.0), Inches(5.0), CARD_BG)
    _add_line(slide, x, Inches(1.8), Inches(4.0), accent, Pt(4))
    # Icon
    _add_circle(slide, x + Inches(0.3), Inches(2.15), Inches(0.65), accent)
    _add_textbox(slide, x + Inches(0.3), Inches(2.15), Inches(0.65), Inches(0.65),
                 icon, font_size=22, colour=DEEP_NAVY, alignment=PP_ALIGN.CENTER)
    # Title
    _add_textbox(slide, x + Inches(1.15), Inches(2.25), Inches(2.6), Inches(0.45),
                 title, font_size=17, colour=WHITE, bold=True)

    y = Inches(3.1)
    for item in items:
        _add_textbox(slide, x + Inches(0.3), y, Inches(0.3), Inches(0.35),
                     "▸", font_size=14, colour=accent)
        _add_textbox(slide, x + Inches(0.65), y, Inches(3.0), Inches(0.35),
                     item, font_size=13, colour=SOFT_WHITE)
        y += Inches(0.42)

    x += Inches(4.25)

_slide_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – APP FEATURES DEEP DIVE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
             "KEY APP FEATURES", font_size=36, colour=WHITE, bold=True)

features = [
    ("🗺", "Interactive Map", "Discover nearby shops on a live map with real-time stock indicators and quick navigation.", ACCENT_CYAN),
    ("🛒", "Smart Cart", "Build procurement drafts, adjust quantities, and manage orders from multiple shops.", GREEN),
    ("📱", "Social Feed", "Browse offers, events, and trade intelligence from retailers and wholesalers.", ACCENT_BLUE),
    ("📷", "Camera Listing", "Snap photos to instantly create product listings with image picker integration.", ACCENT_PURPLE),
    ("🔐", "Multi-Auth", "Email OTP, Google Sign-In, and Phone authentication for secure, flexible access.", ORANGE),
    ("🏷", "Trade Codes", "Scan trade codes from retailers to instantly add mapped SKUs to your cart.", RED_ACCENT),
]

positions = [
    (Inches(0.3), Inches(1.8)),
    (Inches(4.5), Inches(1.8)),
    (Inches(8.7), Inches(1.8)),
    (Inches(0.3), Inches(4.5)),
    (Inches(4.5), Inches(4.5)),
    (Inches(8.7), Inches(4.5)),
]

for (icon, title, desc, accent), (px, py) in zip(features, positions):
    _add_icon_card(slide, px, py, Inches(3.9), Inches(2.3), icon, title, desc, accent)

_slide_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – TOOLS & TECHNOLOGIES
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "TOOLS & TECHNOLOGIES", font_size=36, colour=WHITE, bold=True)

tech_stack = [
    ("Frontend", "Flutter & Dart", "Cross-platform mobile framework for building natively compiled apps from a single codebase.", ACCENT_CYAN),
    ("Backend", "Firebase", "Google's Backend-as-a-Service for real-time database, authentication, and cloud functions.", ACCENT_BLUE),
    ("Auth", "Firebase Auth", "Multi-provider authentication: Email OTP, Google Sign-In, Phone verification.", ACCENT_PURPLE),
    ("Database", "Cloud Firestore", "NoSQL document database with real-time sync, offline support, and automatic scaling.", GREEN),
    ("State Mgmt", "BLoC Pattern", "Business Logic Component pattern for predictable, testable state management in Flutter.", ORANGE),
    ("Maps", "Flutter Map + LatLong2", "Open-source map widget with OpenStreetMap tiles for location-based shop discovery.", RED_ACCENT),
]

x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
y_positions = [Inches(1.7), Inches(4.5)]

for idx, (category, name, desc, accent) in enumerate(tech_stack):
    px = x_positions[idx % 3]
    py = y_positions[idx // 3]

    card = _add_rect(slide, px, py, Inches(3.9), Inches(2.4), CARD_BG)
    _add_line(slide, px, py, Inches(3.9), accent, Pt(3))

    _add_textbox(slide, px + Inches(0.2), py + Inches(0.2), Inches(3.5), Inches(0.3),
                 category.upper(), font_size=11, colour=accent, bold=True)
    _add_textbox(slide, px + Inches(0.2), py + Inches(0.55), Inches(3.5), Inches(0.4),
                 name, font_size=20, colour=WHITE, bold=True)
    _add_textbox(slide, px + Inches(0.2), py + Inches(1.05), Inches(3.5), Inches(1.1),
                 desc, font_size=12, colour=LIGHT_GREY)

_slide_number(slide, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – IMPLEMENTATION DETAILS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
             "IMPLEMENTATION DETAILS", font_size=36, colour=WHITE, bold=True)

impl_items = [
    ("Clean Architecture", "Feature-based folder structure with domain, data, and presentation layers for maintainability.", ACCENT_CYAN),
    ("Role-Based Access", "Three user roles (Customer, Retailer, Wholesaler) with distinct UI flows and permissions.", ACCENT_BLUE),
    ("Real-time Sync", "Firestore streams & listeners for live inventory updates and social feed changes.", ACCENT_PURPLE),
    ("Profile Setup Flow", "Guided onboarding with shop location picker, business details, and category selection.", GREEN),
    ("Marketplace Service", "Singleton service managing seeded shop data, product catalogs, and cart operations.", ORANGE),
    ("Navigation System", "IndexedStack with 4 main tabs (Social, Cart, Map, Profile) and top bar with contextual actions.", RED_ACCENT),
]

y = Inches(1.7)
for title, desc, accent in impl_items:
    _add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), CARD_BG)
    _add_line(slide, Inches(0.8), y, Inches(0.03), accent, Inches(0.8))
    _add_textbox(slide, Inches(1.2), y + Inches(0.08), Inches(3.5), Inches(0.35),
                 title, font_size=16, colour=WHITE, bold=True)
    _add_textbox(slide, Inches(1.2), y + Inches(0.4), Inches(10.8), Inches(0.4),
                 desc, font_size=12, colour=LIGHT_GREY)
    y += Inches(0.92)

_slide_number(slide, 13, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CONCLUSION
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "CONCLUSION", font_size=36, colour=WHITE, bold=True)

_add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8),
             "Nutonium successfully demonstrates that a hyperlocal social commerce platform can bridge\nthe gap between offline retail and digital commerce, benefiting all stakeholders.",
             font_size=16, colour=LIGHT_GREY)

conclusions = [
    ("✓", "Successfully integrates Retail, Supplier, and Customer workflows into a single, cohesive platform.", GREEN),
    ("✓", "Improves Inventory Transparency with real-time stock visibility across all connected retail shops.", ACCENT_CYAN),
    ("✓", "Reduces Product Unavailability through preorder systems and demand-driven supply chain communication.", ACCENT_BLUE),
    ("✓", "Enhances Local Shopping Experience with map-based discovery, social feeds, and instant cart management.", ACCENT_PURPLE),
]

y = Inches(2.8)
for marker, text, accent in conclusions:
    card = _add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.8), CARD_BG)
    _add_circle(slide, Inches(1.7), y + Inches(0.15), Inches(0.5), accent)
    _add_textbox(slide, Inches(1.7), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 marker, font_size=18, colour=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(2.4), y + Inches(0.15), Inches(9.0), Inches(0.55),
                 text, font_size=15, colour=SOFT_WHITE)
    y += Inches(1.0)

_slide_number(slide, 14, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – FUTURE ENHANCEMENTS
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)
_section_header_decor(slide)

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "FUTURE ENHANCEMENTS", font_size=36, colour=WHITE, bold=True)

enhancements = [
    ("🤖", "AI-Based Demand Prediction",
     "Machine learning models to predict product demand based on historical sales, seasonality, and local trends – enabling proactive restocking.",
     ACCENT_CYAN),
    ("🚚", "Delivery Partner Collaboration",
     "Integration with local delivery partners for last-mile fulfillment, turning preorders into doorstep deliveries.",
     ACCENT_BLUE),
    ("📢", "Advertisement Platform",
     "Monetization through targeted, hyperlocal advertisements from external brands and businesses within the platform.",
     ACCENT_PURPLE),
    ("⭐", "Smart Recommendation System",
     "Personalized product and shop recommendations powered by user behavior, preferences, and purchase history.",
     ORANGE),
]

x = Inches(0.3)
for icon, title, desc, accent in enhancements:
    card = _add_rect(slide, x, Inches(1.8), Inches(3.0), Inches(5.0), CARD_BG)
    _add_line(slide, x, Inches(1.8), Inches(3.0), accent, Pt(4))

    _add_circle(slide, x + Inches(0.2), Inches(2.15), Inches(0.7), accent)
    _add_textbox(slide, x + Inches(0.2), Inches(2.15), Inches(0.7), Inches(0.7),
                 icon, font_size=24, colour=DEEP_NAVY, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, x + Inches(0.2), Inches(3.05), Inches(2.6), Inches(0.65),
                 title, font_size=16, colour=WHITE, bold=True)
    _add_textbox(slide, x + Inches(0.2), Inches(3.7), Inches(2.6), Inches(2.8),
                 desc, font_size=12, colour=LIGHT_GREY)

    x += Inches(3.2)

_slide_number(slide, 15, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – REFERENCES & THANK YOU
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_gradient_bg(slide)

# Accent bar top
_add_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT_CYAN, Pt(4))

_add_textbox(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.6),
             "REFERENCES", font_size=30, colour=WHITE, bold=True)

refs = [
    "1.  Flutter Team, Flutter Documentation. Google LLC, 2023.",
    "2.  Firebase Team, Firebase Documentation. Google LLC, 2023.",
    "3.  L. Moroney, The Definitive Guide to Firebase. Apress, 2021.",
    "4.  M. S. Yadav & P. A. Pavlou, \"Marketing in computer-mediated environments,\" J. of Marketing, vol. 78, 2020.",
    "5.  A. Kumar & A. Mukherjee, \"Mobile commerce adoption among small retailers,\" IJRDM, vol. 50, 2022.",
    "6.  D. F. Ferraiolo & D. R. Kuhn, \"Role-based access controls for mobile apps,\" IEEE S&P, vol. 18, 2020.",
    "7.  S. Bhattacharya & A. Sinha, \"Hyperlocal e-commerce: Tech architecture,\" J. Bus. Strategy, vol. 42, 2021.",
    "8.  K. L. Hsiao & C. C. Chen, \"Value-based adoption of social commerce,\" JECR, vol. 23, 2022.",
    "9.  D. Grewal et al., \"The future of retailing,\" J. Retailing, vol. 95, 2020.",
    "10. K. L. Nance & B. Hay, \"Real-time database architectures for mobile apps,\" IJMC, vol. 12, 2021.",
]

y = Inches(1.5)
for ref in refs:
    _add_textbox(slide, Inches(0.8), y, Inches(11.5), Inches(0.3),
                 ref, font_size=10, colour=LIGHT_GREY)
    y += Inches(0.32)

# Thank You section
_add_line(slide, Inches(4), Inches(5.2), Inches(5.3), ACCENT_CYAN, Pt(2))

_add_textbox(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(1.0),
             "THANK YOU", font_size=48, colour=ACCENT_CYAN, bold=True,
             alignment=PP_ALIGN.CENTER)

_add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "Abhishek Bruno  ·  Muhammad Aslam  ·  Gokul SR  ·  Rijin Reji",
             font_size=16, colour=SOFT_WHITE, alignment=PP_ALIGN.CENTER)

_slide_number(slide, 16, TOTAL_SLIDES)

# ── Save ────────────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           "Nutonium_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
